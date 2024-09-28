from tkinter import N
from fastapi.responses import FileResponse
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from docx import Document
from pptx import Presentation
from pptx.util import Pt  # Import Pt from pptx.util
import os

app = FastAPI()

# Allow CORS for all origins
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["POST"],  # Restrict to the methods you need
    allow_headers=["*"],  # You can also restrict headers if needed
)

@app.post("/convert")
async def convert(docx_file: UploadFile = File(...)):
    if docx_file.content_type != "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        raise HTTPException(status_code=400, detail="Invalid file type. Only DOCX files are allowed.")
    output_pptx = 'output.pptx'
    doc = Document(docx_file.file)
    prs = Presentation()

    chunks = extract_chunks(doc)
    for title, lines in chunks:
        add_slide_with_text(prs, title, lines)

    prs.save(output_pptx)
    return FileResponse(output_pptx, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="output.pptx")

def extract_chunks(doc):
    """Extract chunks of text from the document."""
    lines = []
    chunks = []
    current_chunk = []
    current_title = ""

    for paragraph in doc.paragraphs:
        if paragraph.style.name.startswith('Heading'):
            if current_chunk:
                chunks.append((current_title, current_chunk))
                current_chunk = []
            current_title = paragraph.text
        else:
            if paragraph.text.strip() != "":
                current_chunk.append(paragraph.text)
                if len(current_chunk) == 4:
                    chunks.append((current_title, current_chunk))
                    current_chunk = []

    if current_chunk:
        chunks.append((current_title, current_chunk))

    return chunks

def add_slide_with_text(prs, title_text, lines):
    """Add a slide with the given title and lines of text."""
    slide_layout = prs.slide_layouts[1]  # Use the "Title and Content" layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    content.text = "\n".join(lines)

    # Set title font size
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(60)

    remove_bullets(content.text_frame)
    set_white_text_formatting(content.text_frame)
    set_slide_background(slide, RGBColor(0, 0, 0))

def remove_bullets(text_frame):
    """Remove bullets from the text frame."""
    for paragraph in text_frame.paragraphs:
        pPr = paragraph._element.get_or_add_pPr()
        if pPr is not None:
            # Remove bullet properties
            buNone = pPr.find("a:buNone", namespaces=pPr.nsmap)
            if buNone is None:
                buNone = parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
                pPr.append(buNone)
def set_white_text_formatting(text_frame):
    """Set text formatting to white with a specific font size."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = Pt(54)

def set_slide_background(slide, color):
    """Set the slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color