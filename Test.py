from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load the uploaded presentation
pptx_file = r"C:\Users\manoj\Downloads\songs_presentation_dark.pptx"
prs = Presentation(pptx_file)

# Dark background color and white text formatting
def set_white_text_formatting(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text
            run.font.size = Pt(24)  # Set font size to 24pt
        paragraph.alignment = PP_ALIGN.CENTER  # Center the text

# Remove bullets and center text for all slides
def remove_bullets_and_center(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            remove_bullets(shape.text_frame)
            set_white_text_formatting(shape.text_frame)

def remove_bullets(text_frame):
    for paragraph in text_frame.paragraphs:
        paragraph.level = 0  # Ensure no bullet levels are set
        paragraph.text = paragraph.text.lstrip('•\u2022')  # Remove bullet characters if present
        pPr = paragraph._element.get_or_add_pPr()  # Get or add paragraph properties
        if pPr is not None:
            pPr.clear()  # Clear bullet elements

# Apply to all slides
for slide in prs.slides:
    remove_bullets_and_center(slide)

# Save the updated presentation
output_pptx = r"C:\Users\manoj\Downloads\songs_presentation__no_bullets_centered.pptx"
prs.save(output_pptx)

print(f"Updated presentation saved to {output_pptx}")