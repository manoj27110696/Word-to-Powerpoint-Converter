from pptx import Presentation
from docx import Document
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
import os

# Load the word document -add the path to the location of the word document
docx_file = r"C:\Users\manoj\Downloads\CAFC - Sep2024.docx"
doc = Document(docx_file)

# Create a new presentation
prs = Presentation()

# Function to set white text formatting and center the text
def set_white_text_formatting(text_frame, font_size):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text
            run.font.size = Pt(font_size)  # Set font size
        paragraph.alignment = PP_ALIGN.CENTER  # Center the text

# Function to remove bullets
def remove_bullets(text_frame):
    for paragraph in text_frame.paragraphs:
        paragraph.level = 0  # Ensure no bullet levels are set
        pPr = paragraph._element.get_or_add_pPr()  # Get or add paragraph properties
        if pPr is not None:
            buNone = pPr.find('.//a:buNone', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if buNone is None:
                buNone = parse_xml('<a:buNone %s/>' % nsdecls('a'))
                pPr.append(buNone)

# Function to set the background color of a slide
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Function to add a slide with text
def add_slide_with_text(prs, title_text, lines):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    content.text = "\n".join(lines)

    # Set font sizes
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(60)  # Set header font size to 60

    # Remove bullets and center text
    remove_bullets(content.text_frame)
    set_white_text_formatting(content.text_frame, 54)  # Set text font size to 54

    # Set the background color to dark
    set_slide_background(slide, RGBColor(0, 0, 0))  # Dark background

# Split the text into chunks of 4 lines, keeping headers
lines = []
chunks = []
current_chunk = []
current_title = ""

for paragraph in doc.paragraphs:
    if paragraph.style.name.startswith('Heading'):
        if current_chunk:
            chunks.append((current_title, current_chunk))
            current_chunk = []
        current_title = paragraph.text
    else:
        if paragraph.text.strip() != "":
            current_chunk.append(paragraph.text)
            if len(current_chunk) == 4:
                chunks.append((current_title, current_chunk))
                current_chunk = []

if current_chunk:
    chunks.append((current_title, current_chunk))

# Add each chunk as a new slide
for title, chunk in chunks:
    add_slide_with_text(prs, title, chunk)

# Save the updated presentation -add  the path for the location to save the presentation
output_pptx = r"C:\Users\manoj\Downloads\songs_presentation__no_bullets_centered.pptx"
prs.save(output_pptx)

# Check if the output file exists
if os.path.exists(output_pptx):
    print("Output file exists")
else:
    print("Output file does not exist")

print(f"Updated presentation saved to {output_pptx}")